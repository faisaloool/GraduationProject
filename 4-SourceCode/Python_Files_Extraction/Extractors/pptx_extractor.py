from pptx import Presentation

def extract_text_from_pptx(file_path):
    """Extract all text from a PPTX file."""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        return f"Error: {e}"
    return text
