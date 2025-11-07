import os
import comtypes.client
from pptx import Presentation

def extract_text_from_pptx(file_path):
    """Extract all text from a PPTX file."""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        return f"Error: {e}"
    return text

def extract_text_from_ppt(file_path):
    """Convert PPT to temporary PPTX, extract text, then delete temp file."""
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    #powerpoint.Visible = 0

    file_path = os.path.abspath(file_path)
    presentation = powerpoint.Presentations.Open(file_path)
    temp_pptx = file_path + "_temp.pptx"
    presentation.SaveAs(temp_pptx, 24)  # 24 = pptx format
    presentation.Close()
    powerpoint.Quit()

    text = extract_text_from_pptx(temp_pptx)
    os.remove(temp_pptx)
    return text